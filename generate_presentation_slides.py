from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_presentation():
    prs = Presentation()

    # Colors and styles
    green_color = RGBColor(102, 204, 0)
    title_font_size = Pt(40)
    content_font_size = Pt(20)

    # Helper function for adding content slides
    def add_content_slide(title, content_points):
        layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()
        for point in content_points:
            p = text_frame.add_paragraph()
            p.text = point
            p.font.size = content_font_size
            p.space_after = Pt(10)

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "AI-Powered Multi-Cloud\nPenetration Testing Automation"
    slide.placeholders[1].text = "Your Name\nIntroduction to Artificial Intelligence (DA599A)\nKristianstad University, Spring 2025"

    # Slide 2: Interest in AI
    add_content_slide("My Interest in AI", [
        "Edit this.",
        "Edit this.",
        "Edit this."
    ])

    # Slide 3: Motivation Behind Environment Choice
    add_content_slide("Motivation for Multi-Cloud Environment", [
        "Edit this.",
        "Edit this.",
        "Edit this"
    ])

    # Slide 4: Architecture Overview
    layout = prs.slide_layouts[5]  # Title Only layout
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "AI Agent Architecture Overview"
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
    tf = txBox.text_frame
    tf.text = "AI Pentesting Manager\n    ↳ Reconnaissance Agent\n    ↳ Exploitation Agent\n    ↳ Privilege Escalation Agent\n    ↳ Reporting Agent"
    tf.paragraphs[0].font.size = Pt(24)

    # Slide 5: Goal, Percepts, Actions, and Rules
    add_content_slide("Goals, Percepts, Actions & Inference", [
        "Goal: Automate penetration testing intelligently.",
        "Percepts: Cloud configuration data, vulnerabilities, exploits.",
        "Actions: Scanning, exploiting, privilege escalation, reporting.",
        "Inference: Deductive, inductive, abductive, case-based reasoning."
    ])

    # Slide 6: Spotlight on Reconnaissance Agent
    add_content_slide("Reconnaissance Agent Deep Dive", [
        "Edit this.",
        "Edit this.",
        "Edit this."
    ])

    # Slide 7: Additional AI Tooling
    add_content_slide("Integration of AI Tools", [
        "Edit this.",
        "Edit this.",
        "Edit this."
    ])

    # Slide 8: Document Snapshot and References
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Project Documentation & References"
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
    tf = txBox.text_frame
    tf.text = "Insert screenshot here."

    # Slide 9: Code Illustration
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Code Illustration & Snippet"
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
    tf = txBox.text_frame
    tf.text = "Insert Python code snippet."

    # Slide 10: Summary & Future Outlook
    add_content_slide("Summary & Future Outlook", [
        "Edit this.",
        "Edit this.",
        "Edit this."
    ])

    # Style enhancements (colors, alignment, etc.)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.name = 'Arial'
                    paragraph.alignment = PP_ALIGN.LEFT
                    paragraph.font.color.rgb = RGBColor(0, 0, 0)  # black text
                    paragraph.space_after = Pt(6)

        # Title styling
        slide.shapes.title.text_frame.paragraphs[0].font.size = title_font_size
        slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = green_color

    prs.save('AI_Agent_Architecture_Presentation.pptx')


if __name__ == "__main__":
    create_presentation()
    print("Congrats! AI Agent Architecture presentation created successfully!")
